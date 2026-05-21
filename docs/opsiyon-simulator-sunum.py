"""
Opsiyon Dersi Simülatörü — Eğitim Sunumu
PowerPoint dosyası oluşturur: opsiyon-simulator-sunum.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Renk paleti (modern, koyu mavi/yeşil aksanlı)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x1D, 0x9E, 0x75)       # yeşil (kâr)
DANGER = RGBColor(0xD8, 0x5A, 0x30)       # turuncu/kırmızı
INFO = RGBColor(0x3B, 0x82, 0xF6)         # mavi
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF4)
GRAY_TEXT = RGBColor(0x57, 0x53, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, size=18, bold=False, color=NAVY, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_accent_bar(slide, left, top, height=Inches(0.08), width=Inches(1.5), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_card(slide, left, top, width, height, fill=WHITE, border=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if border is None:
        card.line.fill.background()
    else:
        card.line.color.rgb = border
        card.line.width = Pt(0.75)
    return card


# ============================================================
# SLIDE 1 — KAPAK
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.2),
         "Opsiyon Dersi Simülatörü",
         size=54, bold=True, color=WHITE)

add_accent_bar(s, Inches(0.85), Inches(3.4), width=Inches(1.2), color=ACCENT)

add_text(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
         "Greek'lerin opsiyon fiyatına etkisini interaktif olarak öğren",
         size=22, color=WHITE)

add_text(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.4),
         "Yahoo Finance canlı veri  ·  Black-Scholes modeli  ·  Görsel senaryo testleri",
         size=14, color=RGBColor(0xC8, 0xCC, 0xD4))

add_text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
         "Eğitim Sunumu",
         size=12, color=RGBColor(0x9C, 0xA3, 0xAF))


# ============================================================
# SLIDE 2 — AMAÇ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GRAY_LIGHT)

add_accent_bar(s, Inches(0.8), Inches(0.7), color=ACCENT)
add_text(s, Inches(0.8), Inches(0.9), Inches(11.7), Inches(0.8),
         "Amaç", size=40, bold=True, color=NAVY)

add_text(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(0.8),
         "Opsiyon fiyatı neye, ne kadar tepki verir?",
         size=24, bold=True, color=NAVY)

bullets = [
    "Gerçek para riske etmeden, gerçek piyasa verisiyle dene",
    "Delta · Theta · Vega ve IV / spot / zaman etkilerini canlı gör",
    '"Eğer spot şuraya gelirse, IV şu kadar düşerse → primim ne olur?" sorusunun cevabı',
    "Black-Scholes'un kalbini sezgisel olarak kavra",
]
for i, b in enumerate(bullets):
    add_text(s, Inches(1.0), Inches(3.2 + i * 0.7), Inches(0.4), Inches(0.5),
             "▸", size=22, bold=True, color=ACCENT)
    add_text(s, Inches(1.5), Inches(3.2 + i * 0.7), Inches(11.0), Inches(0.5),
             b, size=18, color=NAVY)


# ============================================================
# SLIDE 3 — GENEL ÖZELLİKLER (KART GRİDİ)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)

add_accent_bar(s, Inches(0.8), Inches(0.5), color=ACCENT)
add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7),
         "Genel Özellikler", size=36, bold=True, color=NAVY)

features = [
    ("Canlı Veri", "Yahoo Finance opsiyon zinciri:\nvade, strike, IV, bid/ask", INFO),
    ("İki Eğri Karşılaştırma", "Orijinal Greek (gri) vs.\nslider ile oynanan (yeşil)", ACCENT),
    ("İki Grafik", "Fiyat grafiği (spot aralığı)\nK/Z grafiği (zaman ekseni)", NAVY),
    ("Sliderlar", "Delta · Theta · Vega\n+ spot / IV / kalan gün", DANGER),
    ("Kontrat Adedi", "1 kontrat = 100 hisse\nToplam maliyet otomatik", INFO),
    ("Sözlük + İpuçları", "ITM/ATM/OTM/IV dinamik\naçıklamalar her kutucukta", ACCENT),
]

col_w = Inches(3.9)
row_h = Inches(1.7)
gap = Inches(0.15)
start_left = Inches(0.8)
start_top = Inches(1.8)

for i, (title, desc, accent) in enumerate(features):
    row, col = divmod(i, 3)
    left = start_left + (col_w + gap) * col
    top = start_top + (row_h + gap) * row

    card = add_card(s, left, top, col_w, row_h, fill=GRAY_LIGHT)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             left + Inches(0.2), top + Inches(0.2),
                             Inches(0.08), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(s, left + Inches(0.4), top + Inches(0.18), col_w - Inches(0.5), Inches(0.5),
             title, size=16, bold=True, color=NAVY)
    add_text(s, left + Inches(0.4), top + Inches(0.7), col_w - Inches(0.5), Inches(1.0),
             desc, size=11, color=GRAY_TEXT)


# ============================================================
# SLIDE 4 — TEMEL TERİMLER (Sözlük)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GRAY_LIGHT)

add_accent_bar(s, Inches(0.8), Inches(0.5), color=ACCENT)
add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7),
         "Temel Terimler — Sözlük", size=36, bold=True, color=NAVY)

terms = [
    ("Strike", "Opsiyonu kullandığında hisseyi alacağın (Call) /\nsatacağın (Put) fiyat", INFO),
    ("Prim", "1 hisse başına ödenen opsiyon fiyatı.\n1 kontrat maliyeti = Prim × 100", ACCENT),
    ("ITM", "In The Money: opsiyon kârda\nCall: spot > strike  /  Put: spot < strike", ACCENT),
    ("ATM", "At The Money: strike ≈ spot\n(±%1 yakınında)", INFO),
    ("OTM", "Out of The Money: kârsız taraf\n(Call: spot < strike)", DANGER),
    ("IV", "Implied Volatility — örtük yıllık oynaklık\nYüksek IV = prim pahalı", NAVY),
    ("Delta (Δ)", "Spot 1$ artarsa prim ne kadar değişir\nCall: 0..1, Put: -1..0", ACCENT),
    ("Theta (Θ)", "Günlük zaman erimesi (negatif)\nOTM'de daha hızlı", DANGER),
    ("Vega (ν)", "IV 1 puan artarsa prim ne kadar değişir\nATM'de en yüksek", INFO),
    ("Breakeven", "Vade sonunda kâra geçmek için\nspot bu seviyeyi geçmeli", NAVY),
]

col_w = Inches(3.9)
row_h = Inches(1.0)
gap_y = Inches(0.12)
gap_x = Inches(0.15)
start_left = Inches(0.8)
start_top = Inches(1.7)

for i, (term, desc, accent) in enumerate(terms):
    row, col = divmod(i, 3)
    left = start_left + (col_w + gap_x) * col
    top = start_top + (row_h + gap_y) * row

    card = add_card(s, left, top, col_w, row_h, fill=WHITE)

    add_text(s, left + Inches(0.25), top + Inches(0.1), Inches(1.5), Inches(0.45),
             term, size=15, bold=True, color=accent)
    add_text(s, left + Inches(0.25), top + Inches(0.42), col_w - Inches(0.4), Inches(0.65),
             desc, size=10, color=GRAY_TEXT)


# ============================================================
# SLIDE 5 — KULLANIM AKIŞI (4 ADIM)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)

add_accent_bar(s, Inches(0.8), Inches(0.5), color=ACCENT)
add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7),
         "Kullanım Akışı", size=36, bold=True, color=NAVY)

steps = [
    ("1", "Sembol Ara", "TSLA, AAPL, NVDA...\nşirket adı da yazabilirsin", INFO),
    ("2", "Kontrat Seç", "Vade · Call/Put · Strike\n(ATM yakın 30/100/200)", ACCENT),
    ("3", "Nokta İşaretle", "Grafik alanında mouse ile tıklayarak\ntarihlere göre merak ettiğin\nspot değerlerini işaretle", DANGER),
    ("4", "Slider'larla Oyna", "Spot, IV, gün veya Delta/Theta/Vega\ndoğrudan oynatarak K/Z'yi gör", NAVY),
]

col_w = Inches(2.85)
row_h = Inches(3.5)
gap = Inches(0.25)
start_left = Inches(0.85)
start_top = Inches(2.0)

for i, (num, title, desc, accent) in enumerate(steps):
    left = start_left + (col_w + gap) * i
    top = start_top

    card = add_card(s, left, top, col_w, row_h, fill=GRAY_LIGHT)

    circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                left + Inches(1.0), top + Inches(0.35),
                                Inches(0.85), Inches(0.85))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()

    add_text(s, left + Inches(1.0), top + Inches(0.4), Inches(0.85), Inches(0.75),
             num, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, left + Inches(0.2), top + Inches(1.5), col_w - Inches(0.4), Inches(0.6),
             title, size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_text(s, left + Inches(0.2), top + Inches(2.2), col_w - Inches(0.4), Inches(1.2),
             desc, size=12, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — KONTROL LİSTESİ: Greek Davranışları (Tablo)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, GRAY_LIGHT)

add_accent_bar(s, Inches(0.8), Inches(0.5), color=ACCENT)
add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7),
         "Kontrol Listesi — Greek Davranışları (Call için)",
         size=30, bold=True, color=NAVY)

# Tablo
rows_data = [
    ["Slider", "Beklenen Davranış", "Kontrol Et"],
    ["Spot ↑", "Prim artar (Delta pozitif)", "Eğri sola → yukarı kayıyor mu?"],
    ["Kalan gün ↓", "Prim azalır (Theta erimesi)", "OTM'de daha mı hızlı eriyor?"],
    ["IV ↑", "Prim artar (Vega pozitif)", "ATM'de etki en büyük mü?"],
    ["Spot = Strike", "Delta ≈ 0.5", "Doğrulanıyor mu?"],
    ["Vade günü, spot = strike", "Zaman değeri 0 → tüm prim kayıp", "Doğrulanıyor mu?"],
]

table_left = Inches(0.8)
table_top = Inches(1.9)
table_width = Inches(11.7)
table_height = Inches(4.5)

tbl_shape = s.shapes.add_table(len(rows_data), 3, table_left, table_top, table_width, table_height)
table = tbl_shape.table

# Sütun genişlikleri
table.columns[0].width = Inches(2.8)
table.columns[1].width = Inches(4.4)
table.columns[2].width = Inches(4.5)

for r_idx, row_data in enumerate(rows_data):
    for c_idx, val in enumerate(row_data):
        cell = table.cell(r_idx, c_idx)
        cell.text = ""  # önce temizle
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = val
        run.font.name = 'Calibri'

        if r_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            run.font.color.rgb = WHITE
            run.font.bold = True
            run.font.size = Pt(15)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 1 else GRAY_LIGHT
            run.font.color.rgb = NAVY if c_idx == 0 else GRAY_TEXT
            run.font.bold = (c_idx == 0)
            run.font.size = Pt(13)


# ============================================================
# SLIDE 7 — KRİTİK SORULAR (Kendine sor)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, WHITE)

add_accent_bar(s, Inches(0.8), Inches(0.5), color=DANGER)
add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.7),
         "Kritik Sorular — Kendine Sor", size=36, bold=True, color=NAVY)

questions = [
    ("Theta tuzağı",
     "Aynı spot, sadece 7 gün geçti — kaybın ne kadar?\nÖzellikle OTM'de şok yaşayabilirsin."),
    ("IV crush (IV çöküşü)",
     "IV %58'den %35'e düşerse primim ne olur?\nKazanç sayfasından sonra yaşanır — pozisyonu mahvedebilir."),
    ("Breakeven mantığı",
     "Vade günü spot = strike → kâr mı zarar mı?\nCevap: tam primi kaybedersin (zaman değeri 0)."),
    ("Kontrat adedi etkisi",
     "1 yerine 5 kontrat → dolar K/Z × 5, ama % K/Z aynı.\nOran toplam maliyet bazlı hesaplanır."),
]

for i, (q, a) in enumerate(questions):
    row, col = divmod(i, 2)
    left = Inches(0.8) + (Inches(5.95) + Inches(0.15)) * col
    top = Inches(1.8) + (Inches(2.5) + Inches(0.2)) * row

    card = add_card(s, left, top, Inches(5.95), Inches(2.5), fill=GRAY_LIGHT)

    icon_circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     left + Inches(0.3), top + Inches(0.3),
                                     Inches(0.5), Inches(0.5))
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = DANGER
    icon_circle.line.fill.background()
    add_text(s, left + Inches(0.3), top + Inches(0.32), Inches(0.5), Inches(0.5),
             "?", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, left + Inches(1.0), top + Inches(0.35), Inches(4.8), Inches(0.5),
             q, size=17, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(1.05), Inches(5.4), Inches(1.4),
             a, size=12, color=GRAY_TEXT)


# ============================================================
# SLIDE 8 — UYARI + KAPANIŞ
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

add_accent_bar(s, Inches(0.8), Inches(0.5), color=DANGER, width=Inches(1.5))
add_text(s, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8),
         "⚠ Önemli Uyarı", size=36, bold=True, color=WHITE)

warn_card = add_card(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.5),
                     fill=RGBColor(0x1A, 0x2E, 0x52))

add_text(s, Inches(1.1), Inches(2.2), Inches(11.0), Inches(0.5),
         "Bu araç yalnızca eğitim amaçlıdır.",
         size=22, bold=True, color=DANGER)
add_text(s, Inches(1.1), Inches(2.8), Inches(11.0), Inches(1.7),
         "• Gerçek alım-satım kararı için kullanma\n"
         "• Black-Scholes basitleştirilmiş bir modeldir (kâr payı, faiz, Amerikan tipi erken kullanım hariç)\n"
         "• Yahoo'nun IV / bid-ask verisi zaman zaman eksik veya gecikmeli olabilir\n"
         "• Düşük likiditeli kontratlarda fiyat sapmaları büyük olur",
         size=14, color=WHITE)

add_text(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.6),
         "Hazırsan başla → bir sembol ara, oyna, sor.",
         size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
         "Opsiyonların kalbi: zaman, oynaklık, ve mesafe.",
         size=14, color=RGBColor(0xC8, 0xCC, 0xD4), align=PP_ALIGN.CENTER)


# ============================================================
# KAYDET
# ============================================================
output = "docs/opsiyon-simulator-sunum.pptx"
prs.save(output)
print(f"OK -> {output}")
print(f"Slayt sayisi: {len(prs.slides)}")
